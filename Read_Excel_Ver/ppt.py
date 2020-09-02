"""
Function: generate a ppt file based on the figures' directory and the given abnormal values' list
Author: Xinran Wang
Date: 08/14/2020
"""

import os
from os import listdir
from os.path import join, basename
from re import findall
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from change_line_test import *


def generate_ppt(fig_path, abnormal_dic, target_dir):
    """
    Generate the powerpoint, each page has one figure; for the merged dataframe, the outlier interval will also be listed on the powerpoint; the powerpoint will be saved in the same directory as this project
    :param fig_path: the absolute path of the folder containing all the figures drawn
    :param abnormal_dic: a dictionary with keys as the file names of data figures that needed to be plotted with abnormals, values as the corresponding outlier values
    :param target_dir: the target directory
    :return: None
    """
    # set as reading png file, subject to change (?)
    pic_files = [join(fig_path, fn) for fn in listdir(fig_path) if fn.endswith(".png")]

    ppt_file = pptx.Presentation()
    ppt_file.slide_width = Inches(16)
    ppt_file.slide_height = Inches(9)
    for fn in pic_files:
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        txt = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), ppt_file.slide_width, Inches(5))
        p = txt.text_frame.add_paragraph()
        slide.shapes.add_picture(fn, Inches(0), Inches(1.5), Inches(16), Inches(6.4))

        pic_name = basename(fn)
        splitted_name = pic_name.split("-")
        data_type = splitted_name[0]
        p.text = data_type + ": data comparison figure"
        p.font.bold = True

        if "Stats" in pic_name:
            abnormal_text, lines = change_lines(abnormal_dic[data_type])
            if lines <= 1:
                abnormal = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), ppt_file.slide_width, Inches(5))
                para = abnormal.text_frame.add_paragraph()
                para.text = "Potential abnormal data Camera ID ranges:" + "\n" + abnormal_text
            else:
                slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
                abnormal = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), ppt_file.slide_width, Inches(5))
                para = abnormal.text_frame.add_paragraph()
                para.text = "Potential abnormal data Camera IDs:" + "\n" + abnormal_text

        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = "Times New Roman"
        p.font.size = Pt(25)
    ppt_file.save(os.path.join(target_dir, "Reinjectiontest2.pptx"))
